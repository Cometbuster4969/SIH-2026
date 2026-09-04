"""SatQuery AI — SIH 2026 idea deck builder (PS 26167, ISRO).

One layout spec  ->  two renderers:
  * PPTX  : built ON the official SIH2026-IDEA-Presentation-Format.pptx template
            (chrome, masters, footer, SIH logo kept; instruction slide removed)
  * PDF   : rendered directly with PyMuPDF at identical geometry (SIH requires
            a PDF upload; no LibreOffice available in this environment)

Run:  python3 build_deck.py
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import fitz  # PyMuPDF

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, "SIH2026-IDEA-Presentation-Format.pptx")
ASSETS = os.path.join(HERE, "assets")
OUT_PPTX = os.path.join(ROOT, "SatQuery-AI-SIH2026-Idea-Deck.pptx")
OUT_PDF = os.path.join(ROOT, "SatQuery-AI-SIH2026-Idea-Deck.pdf")

SW, SH = 13.3333, 7.5           # slide inches
PT = 72.0                       # pdf points per inch

# ---------------------------------------------------------------- palette
BLUE = "0070C0"; DBLUE = "1F497D"; ORANGE = "E87722"; GREEN = "2E9E5B"
RED = "D93025"; TEAL = "12B5A5"; PURPLE = "7030A0"; GREY = "58606A"
INK = "101418"; LGREY = "E8EBEF"; TINT_B = "EAF3FB"; TINT_G = "E9F7EE"
TINT_P = "F3EEFA"; TINT_O = "FFF4E8"; WHITE = "FFFFFF"; CARDGREY = "EEF1F4"

F_ARIAL = "Arial"; F_SERIF = "Times New Roman"; F_GAR = "Garamond"
PDF_FONT = {F_ARIAL: ("arimo", "arimob", "arimoi", "arimobi"),
            F_SERIF: ("tinos", "tinosb", "tinosi", "tinosbi"),
            F_GAR: ("tinos", "tinosb", "tinosi", "tinosbi")}
FONTFILES = {
    "arimo": "/usr/share/fonts/truetype/Arimo-Regular.ttf",
    "arimob": "/usr/share/fonts/truetype/Arimo-Bold.ttf",
    "arimoi": "/usr/share/fonts/truetype/Arimo-Italic.ttf",
    "arimobi": "/usr/share/fonts/truetype/Arimo-BoldItalic.ttf",
    "tinos": "/usr/share/fonts/truetype/Tinos-Regular.ttf",
    "tinosb": "/usr/share/fonts/truetype/Tinos-Bold.ttf",
    "tinosi": "/usr/share/fonts/truetype/Tinos-Italic.ttf",
    "tinosbi": "/usr/share/fonts/truetype/Tinos-BoldItalic.ttf",
}
_FCACHE = {}


def hx(s):
    s = s.lstrip("#")
    return tuple(int(s[i:i + 2], 16) / 255.0 for i in (0, 2, 4))


def ffont(key):
    if key not in _FCACHE:
        _FCACHE[key] = fitz.Font(fontfile=FONTFILES[key])
    return _FCACHE[key]


# ---------------------------------------------------------------- spec model
class Run:
    def __init__(self, text, size=12, bold=False, italic=False, color=INK, font=F_ARIAL):
        self.text, self.size, self.bold, self.italic, self.color, self.font = \
            text, size, bold, italic, color, font


class Para:
    def __init__(self, runs, align="left", space_before=0, space_after=4, line=1.12,
                 bullet=None):
        self.runs = runs if isinstance(runs, list) else [runs]
        self.align, self.space_before, self.space_after, self.line = \
            align, space_before, space_after, line
        self.bullet = bullet          # glyph prefix drawn by renderer


class El:
    """kind: rect | rrect | oval | line | arrow | text | image | hexagon"""
    def __init__(self, kind, **kw):
        self.kind = kind
        self.__dict__.update(kw)


def rect(x, y, w, h, fill=None, line=None, lw=1.0, radius=None, dashed=False):
    kind = "rrect" if radius else "rect"
    return El(kind, x=x, y=y, w=w, h=h, fill=fill, line=line, lw=lw,
              radius=radius or 0, dashed=dashed)


def oval(x, y, w, h, fill=None, line=None, lw=1.0):
    return El("oval", x=x, y=y, w=w, h=h, fill=fill, line=line, lw=lw)


def line(x1, y1, x2, y2, color=GREY, lw=1.2, dashed=False, arrow=False):
    return El("arrow" if arrow else "line", x1=x1, y1=y1, x2=x2, y2=y2,
              color=color, lw=lw, dashed=dashed)


def text(x, y, w, h, paras, anchor="top", pdf_only=False):
    return El("text", x=x, y=y, w=w, h=h, paras=paras, anchor=anchor,
              pdf_only=pdf_only)


def image(path, x, y, w, h):
    return El("image", path=path, x=x, y=y, w=w, h=h)


def P(txt, size=12, bold=False, italic=False, color=INK, font=F_ARIAL,
      align="left", space_after=4, space_before=0, line=1.12):
    return Para([Run(txt, size, bold, italic, color, font)], align,
                space_before, space_after, line)


def rich(pairs, align="left", space_after=4, space_before=0, line=1.12, bullet=None):
    """pairs: list of (text, size, bold, italic, color[, font]) tuples"""
    runs = []
    for t in pairs:
        txt, size, bold, italic, color = t[0], t[1], t[2], t[3], t[4]
        font = t[5] if len(t) > 5 else F_ARIAL
        runs.append(Run(txt, size, bold, italic, color, font))
    return Para(runs, align, space_before, space_after, line, bullet)


# ---------------------------------------------------------------- PPTX backend
def _set_font(run, r):
    f = run.font
    f.size = Pt(r.size)
    f.bold = r.bold
    f.italic = r.italic
    f.name = r.font
    f.color.rgb = RGBColor.from_string(r.color.lstrip("#"))


def _add_text(slide, el):
    tb = slide.shapes.add_textbox(Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[el.anchor]
    for i, para in enumerate(el.paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[para.align]
        p.space_before = Pt(para.space_before)
        p.space_after = Pt(para.space_after)
        p.line_spacing = para.line
        runs = list(para.runs)
        if para.bullet:
            r0 = p.add_run(); r0.text = para.bullet + "  "
            _set_font(r0, runs[0])
        for r in runs:
            rn = p.add_run(); rn.text = r.text
            _set_font(rn, r)
    return tb


def _fill(shape, color):
    if color:
        color = color.lstrip("#")
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(color)
    else:
        shape.fill.background()


def _linefmt(shape, color, lw, dashed=False):
    if color:
        color = color.lstrip("#")
        shape.line.color.rgb = RGBColor.from_string(color)
        shape.line.width = Pt(lw)
        if dashed:
            shape.line.dash_style = 4  # dash
    else:
        shape.line.fill.background()


def _add_shape(slide, el):
    m = {"rect": MSO_SHAPE.RECTANGLE, "rrect": MSO_SHAPE.ROUNDED_RECTANGLE,
         "oval": MSO_SHAPE.OVAL}[el.kind]
    sp = slide.shapes.add_shape(m, Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
    if el.kind == "rrect":
        try:
            sp.adjustments[0] = min(0.5, el.radius / min(el.w, el.h))
        except Exception:
            pass
    _fill(sp, el.fill)
    _linefmt(sp, el.line, el.lw, getattr(el, "dashed", False))
    sp.shadow.inherit = False
    sp.text_frame.text = ""
    return sp


def _add_line(slide, el):
    cn = slide.shapes.add_connector(1, Inches(el.x1), Inches(el.y1),
                                    Inches(el.x2), Inches(el.y2))
    cn.line.color.rgb = RGBColor.from_string(el.color)
    cn.line.width = Pt(el.lw)
    if el.dashed:
        cn.line.dash_style = 4
    if el.kind == "arrow":
        ln = cn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return cn


def render_pptx(slides_spec, out):
    prs = Presentation(TEMPLATE)
    # drop the instruction slide (#7)
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    if len(ids) == 7:
        xml_slides.remove(ids[6])
    slides = list(prs.slides)
    for idx, spec in enumerate(slides_spec):
        slide = slides[idx]
        # clear template placeholder body text we replace ourselves
        for sh in list(slide.shapes):
            if sh.name in ("TextBox 8",):
                sh._element.getparent().remove(sh._element)
        for el in spec:
            if getattr(el, "pdf_only", False):
                continue
            if el.kind == "text":
                _add_text(slide, el)
            elif el.kind in ("rect", "rrect", "oval"):
                _add_shape(slide, el)
            elif el.kind in ("line", "arrow"):
                _add_line(slide, el)
            elif el.kind == "image":
                slide.shapes.add_picture(el.path, Inches(el.x), Inches(el.y),
                                         Inches(el.w), Inches(el.h))
        _apply_slide_texts(prs, slide, idx)
    prs.save(out)
    return out


def _set_ph(slide, name, value, size=None, bold=None, color=None, italic=None,
            font=None, align=None):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            tf = sh.text_frame
            tf.clear()
            lines = value.split("\n")
            for i, ln in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run(); r.text = ln
                if size: r.font.size = Pt(size)
                if bold is not None: r.font.bold = bold
                if italic is not None: r.font.italic = italic
                r.font.name = font or F_ARIAL
                r.font.color.rgb = RGBColor.from_string(color or INK)
                if align:
                    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                                   "right": PP_ALIGN.RIGHT}[align]
            return sh
    return None


# ---------------------------------------------------------------- PDF backend
class Pdf:
    def __init__(self):
        self.doc = fitz.open()

    def page(self):
        pg = self.doc.new_page(width=SW * PT, height=SH * PT)
        for name, path in FONTFILES.items():
            pg.insert_font(fontname=name, fontfile=path)
        return pg

    def R(self, v):
        return v * PT

    def draw(self, pg, el):
        k = el.kind
        if k == "image":
            pg.insert_image(fitz.Rect(self.R(el.x), self.R(el.y),
                                      self.R(el.x + el.w), self.R(el.y + el.h)),
                            filename=el.path)
            return
        if k == "hexagon":
            fill = hx(el.fill) if el.fill else None
            col = hx(el.line) if el.line else None
            sh = pg.new_shape()
            sh.draw_polyline(list(el.pts) + [el.pts[0]])
            sh.finish(color=col, fill=fill, width=el.lw)
            sh.commit()
            return
        if k in ("rect", "rrect", "oval"):
            r = fitz.Rect(self.R(el.x), self.R(el.y),
                          self.R(el.x + el.w), self.R(el.y + el.h))
            fill = hx(el.fill) if el.fill else None
            col = hx(el.line) if el.line else None
            dashes = "[3 2] 0" if getattr(el, "dashed", False) else None
            if k == "oval":
                import math as _m
                cx, cy = (r.x0 + r.x1) / 2, (r.y0 + r.y1) / 2
                rx, ry = (r.x1 - r.x0) / 2, (r.y1 - r.y0) / 2
                pts = [fitz.Point(cx + rx * _m.cos(2 * _m.pi * i / 40),
                                  cy + ry * _m.sin(2 * _m.pi * i / 40))
                       for i in range(40)]
                sh = pg.new_shape()
                sh.draw_polyline(pts + [pts[0]])
                sh.finish(color=col, fill=fill, width=el.lw, dashes=dashes)
                sh.commit()
                return
            elif k == "rrect":
                rad = min(el.radius, min(el.w, el.h) / 2) * PT
                pg.draw_rect(r, color=col, fill=fill, width=el.lw, dashes=dashes,
                             radius=min(0.5, el.radius / min(el.w, el.h)))
            else:
                pg.draw_rect(r, color=col, fill=fill, width=el.lw, dashes=dashes)
            return
        if k in ("line", "arrow"):
            p0 = fitz.Point(self.R(el.x1), self.R(el.y1))
            p1 = fitz.Point(self.R(el.x2), self.R(el.y2))
            dashes = "[3 2] 0" if el.dashed else None
            pg.draw_line(p0, p1, color=hx(el.color), width=el.lw,
                         dashes=dashes)
            if k == "arrow":
                import math
                ang = math.atan2(p1.y - p0.y, p1.x - p0.x)
                L = 6.5
                for da in (2.6, -2.6):
                    pg.draw_line(p1, fitz.Point(p1.x - L * math.cos(ang + da * 0.14),
                                                p1.y - L * math.sin(ang + da * 0.14)),
                                 color=hx(el.color), width=el.lw)
            return
        if k == "text":
            self.draw_text(pg, el)

    def draw_text(self, pg, el):
        y = self.R(el.y)
        x0 = self.R(el.x)
        w = self.R(el.w)
        if el.anchor != "top":
            total = self.measure(el)
            if el.anchor == "middle":
                y = self.R(el.y) + (self.R(el.h) - total) / 2
            else:
                y = self.R(el.y) + self.R(el.h) - total
        for para in el.paras:
            y += para.space_before * PT / 72.0 * 1.0
            lines = self.wrap(para, w)
            lh = max(r.size for r in para.runs) * para.line * PT / 72.0
            for ln in lines:
                # line width for alignment
                lw_ = sum(ffont(self._key(r)).text_length(t, r.size)
                          for (t, r) in ln)
                if para.align == "center":
                    cx = x0 + (w - lw_) / 2
                elif para.align == "right":
                    cx = x0 + w - lw_
                else:
                    cx = x0
                for (t, r) in ln:
                    key = self._key(r)
                    pg.insert_text(fitz.Point(cx, y + lh * 0.80), t,
                                   fontname=key, fontsize=r.size,
                                   color=hx(r.color))
                    cx += ffont(key).text_length(t, r.size)
                y += lh
            y += para.space_after * PT / 72.0
        return y

    def _key(self, r):
        fam = PDF_FONT.get(r.font, PDF_FONT[F_ARIAL])
        if r.bold and r.italic:
            return fam[3]
        return fam[1] if r.bold else (fam[2] if r.italic else fam[0])

    def wrap(self, para, w_pt):
        items = []
        if para.bullet:
            items.append((para.bullet + "   ", para.runs[0]))
        for r in para.runs:
            items.extend((tok, r) for tok in self._tokens(r.text))
        lines, cur, curw = [], [], 0.0
        space_w = {r.size: ffont(self._key(r)).text_length(" ", r.size) for r in para.runs}
        for tok, r in items:
            fw = ffont(self._key(r)).text_length(tok, r.size)
            sw = space_w[r.size]
            if cur and curw + fw > w_pt:
                lines.append(cur); cur, curw = [], 0.0
                if tok == " ":
                    continue
            cur.append((tok, r)); curw += fw + (sw if tok != " " else 0)
        if cur:
            lines.append(cur)
        return lines or [[]]

    def _tokens(self, s):
        out, buf = [], ""
        for ch in s:
            buf += ch
            if ch == " ":
                out.append(buf); buf = ""
        if buf:
            out.append(buf)
        return out

    def measure(self, el):
        total = 0.0
        w_pt = self.R(el.w)
        for para in el.paras:
            lh = max(r.size for r in para.runs) * para.line * PT / 72.0
            total += para.space_before * PT / 72.0
            total += len(self.wrap(para, w_pt)) * lh
            total += para.space_after * PT / 72.0
        return total

    def save(self, out):
        try:
            self.doc.subset_fonts(fallback=False)
        except Exception as e:
            print("subset_fonts:", e)
        self.doc.save(out, garbage=4, deflate=True)
        self.doc.close()


def render_pdf(slides_spec, chrome, out):
    pdf = Pdf()
    for idx, spec in enumerate(slides_spec):
        pg = pdf.page()
        ch = chrome[idx]
        # background
        pg.draw_rect(fitz.Rect(0, 0, SW * PT, SH * PT), fill=hx(WHITE), color=None)
        for el in ch["bg"]:
            pdf.draw(pg, el)
        for el in spec:
            pdf.draw(pg, el)
        for el in ch["fg"]:
            pdf.draw(pg, el)
    pdf.save(out)
    return out


# ================================================================ CONTENT
def chrome_for(idx, title_lines, pointer, number):
    """Template chrome replicated for the PDF renderer + PPTX placeholder edits."""
    bg, fg = [], []
    if idx == 0:
        import math
        def hexagon(cx, cy, r, fill=None, line=None, lw=1.0):
            pts = [fitz.Point(*[0, 0]) for _ in range(6)]
            pp = []
            for k in range(6):
                a = math.radians(60 * k - 90)
                pp.append(fitz.Point((cx + r * math.cos(a)) * PT,
                                     (cy + r * math.sin(a)) * PT))
            pgdraw_hex.append((pp, fill, line, lw))
        pgdraw_hex = []
        hexagon(9.10, 4.05, 2.55, fill=LGREY)
        hexagon(7.85, 1.85, 0.95, line="#DDE1E6", lw=2.2)
        hexagon(6.70, 4.80, 0.55, fill=LGREY)
        for pp, fill, ln, lw in pgdraw_hex:
            bg.append(El("hexagon", pts=pp, fill=fill, line=ln, lw=lw))
        bg.append(image(os.path.join(ASSETS, "sih_bulb.png"), 7.50, 1.88, 3.50, 3.75))
        bg.append(image(os.path.join(ASSETS, "sih_logo.png"), 10.70, 0.0, 2.46, 1.16))
        fg.append(text(0.36, -0.58, 11.33, 2.27,
                       [P("SMART INDIA HACKATHON 2026", 40, True, color=DBLUE,
                          font=F_GAR, space_after=0)], anchor="middle"))
        fg.append(text(0.55, 1.32, 8.0, 0.5,
                       [P("SatQuery AI — ask your satellite in plain language",
                          17, False, True, GREY, font=F_SERIF, space_after=0)]))
        return {"bg": bg, "fg": fg}
    # footer bar
    bg.append(rect(0, 6.95, 13.3333, 0.55, fill=BLUE))
    fg.append(text(0, 7.03, 13.3333, 0.4,
                   [P("SMART INDIA HACKATHON 2026", 13, True, color=WHITE,
                      font=F_SERIF, align="center", space_after=0)], anchor="middle"))
    fg.append(text(11.6, 7.03, 1.35, 0.4,
                   [P(str(number), 13, True, color=WHITE, font=F_SERIF,
                      align="center", space_after=0)], anchor="middle"))
    # SIH logo top-right (image added in both renderers via spec? keep in bg)
    bg.append(image(os.path.join(ASSETS, "sih_logo.png"), 10.70, 0.0, 2.46, 1.16))
    # team badge
    bg.append(oval(0.36, 0.28, 1.37, 0.88, fill=WHITE, line=PURPLE, lw=1.2))
    fg.append(text(0.36, 0.28, 1.37, 0.88,
                   [P("x64", 15, True, color=PURPLE, align="center", space_after=0)],
                   anchor="middle"))
    # title
    fg.append(text(1.90, 0.10, 9.40, 0.85,
                   [P(title_lines, 30, True, color=INK, font=F_SERIF,
                      align="center", space_after=0)], anchor="middle"))
    if pointer:
        fg.append(text(1.90, 1.02, 9.40, 0.30,
                       [P(pointer, 9.5, False, True, GREY, align="center",
                          space_after=0)], anchor="middle"))
    return {"bg": bg, "fg": fg}


def slide1():
    """Title page — template art kept; fields filled like the example deck."""
    els = []
    els.append(text(0.40, 2.20, 6.60, 4.95, [
        rich([("Problem Statement ID – ", 18, True, False, INK),
              ("26167", 18, False, False, INK)], space_after=9),
        rich([("Problem Statement Title – ", 18, True, False, INK),
              ("SatQuery AI: An Interactive Vision-Language Assistant for "
               "Multimodal Remote Sensing Image Analysis through Text Queries",
               18, False, False, INK)], space_after=9, line=1.15),
        rich([("Theme – ", 18, True, False, INK),
              ("Space Technology", 18, False, False, INK)], space_after=9),
        rich([("PS Category – ", 18, True, False, INK),
              ("Software", 18, False, False, INK)], space_after=9),
        rich([("Team ID – ", 18, True, False, INK),
              ("[ to be filled ]", 18, False, False, GREY)], space_after=9),
        rich([("Team Name – ", 18, True, False, INK),
              ("x64", 18, False, False, INK)], space_after=9),
    ], pdf_only=True))
    return els


def slide2():
    els = []
    # left: proposed solution banner + box
    els.append(rect(0.30, 1.50, 5.55, 0.50, fill=BLUE, radius=0.08))
    els.append(text(0.30, 1.50, 5.55, 0.50,
                    [P("Proposed Solution / Approach", 16, True, color=WHITE,
                       font=F_SERIF, align="center", space_after=0)], anchor="middle"))
    els.append(rect(0.30, 2.06, 5.55, 4.79, fill=TINT_B, line=BLUE, lw=1.0,
                    radius=0.06))
    bul = dict(space_after=5, line=1.08)
    els.append(text(0.52, 2.22, 5.12, 4.5, [
        rich([("Agentic by design: ", 11.5, True, False, DBLUE),
              ("a natural-language query is planned into a sequence of specialist "
               "tools drawn from a predefined registry (PS P6–P8) — not one "
               "monolithic VLM.", 11.5, False, False, INK)], **bul),
        rich([("Guardrailed execution: ", 11.5, True, False, DBLUE),
              ("schema + parameter whitelist and a deterministic fallback router; "
               "incompatible inputs return a compatibility report — never a "
               "crash.", 11.5, False, False, INK)], **bul),
        rich([("Auditable trace (the graded surface): ", 11.5, True, False, DBLUE),
              ("every run logs task, tools + versions, parameters, input hashes, "
               "latency and confidence (P9–P10).", 11.5, False, False, INK)], **bul),
        rich([("Evidence, not vibes: ", 11.5, True, False, DBLUE),
              ("answers ship with change masks, grounding boxes and caption "
               "overlays, plus calibrated High/Med/Low badges and an optical↔SAR "
               "disagreement flag.", 11.5, False, False, INK)], **bul),
        rich([("One-click audit report: ", 11.5, True, False, DBLUE),
              ("downloadable PDF/HTML report per query — trace + evidence + "
               "confidence, ready for evaluation.", 11.5, False, False, INK)],
             space_after=8, line=1.10),
        rich([("How it addresses the problem: ", 10.5, True, True, ORANGE),
              ("turns expert GIS pipelines into a query box, and makes every "
               "answer checkable — the trust gap the PS itself names.",
               10.5, False, True, GREY)], space_after=6, line=1.10),
        P("Representative PS queries we demo (verbatim):", 10, True, False, DBLUE,
          space_after=3),
        P("1  “Describe the land-cover and major objects visible in this image.”",
          8.4, False, True, GREY, space_after=1, line=1.0),
        P("2  “Highlight the water body referred to in the query.”",
          8.4, False, True, GREY, space_after=1, line=1.0),
        P("3  “What changed between these two dates, and where did the change occur?”",
          8.4, False, True, GREY, space_after=1, line=1.0),
        P("4  “Use the optical and SAR images together to identify built-up and "
          "water-covered regions.”", 8.4, False, True, GREY, space_after=1, line=1.0),
        P("5  “Has the built-up area increased, decreased, or remained unchanged?”",
          8.4, False, True, GREY, space_after=0, line=1.0),
    ]))
    # right: workflow diagram
    els.append(image(os.path.join(ASSETS, "fig_workflow.png"), 6.02, 1.38, 7.02, 2.30))
    # prototype evidence strip
    ey = 3.94
    for i, (f, cap) in enumerate([("overlay_change.png", "change mask"),
                                  ("overlay_ground.png", "grounding boxes"),
                                  ("overlay_caption.png", "caption overlay")]):
        x = 6.55 + i * 2.05
        els.append(image(os.path.join(ASSETS, f), x, ey, 1.12, 1.12))
        els.append(rect(x, ey, 1.12, 1.12, line=GREY, lw=0.8))
        els.append(text(x - 0.25, ey + 1.14, 1.62, 0.22,
                        [P(cap, 8, True, color=GREY, align="center", space_after=0)]))
    els.append(text(6.02, ey - 0.24, 7.02, 0.22,
                    [P("REAL PROTOTYPE OUTPUTS — thin end-to-end build (Window A)",
                       8.5, True, color=DBLUE, align="center", space_after=0)]))
    # innovation cards
    iy = 5.62
    els.append(rect(7.30, iy - 0.30, 4.40, 0.28, fill=GREEN, radius=0.14))
    els.append(text(7.30, iy - 0.30, 4.40, 0.28,
                    [P("Innovation & Uniqueness", 11, True, color=WHITE,
                       font=F_SERIF, align="center", space_after=0)], anchor="middle"))
    cards = [
        ("Trace as a product", "the graded surface is a first-class UI + JSON artifact", BLUE),
        ("Registry + guardrails", "tools declared, params whitelisted, routing deterministic", ORANGE),
        ("Optical × SAR joint reasoning", "optical gives spectra & context; SAR sees through cloud, day & night — paired focus (P1, P5)", GREEN),
        ("RS-adapted, never generic", "QLoRA on BigEarthNet.txt — generic VLMs explicitly insufficient (P12)", PURPLE),
        ("Honest confidence", "calibrated badges, disagreement flags, reasoned abstention", TEAL),
    ]
    cw = 1.365
    for i, (t, b, col) in enumerate(cards):
        x = 6.02 + i * (cw + 0.05)
        els.append(rect(x, iy, cw, 1.22, fill=col, radius=0.07))
        els.append(text(x + 0.05, iy + 0.07, cw - 0.10, 1.10, [
            P(t, 8.6, True, color=WHITE, align="center", space_after=2, line=1.02),
            P(b, 6.9, False, color=WHITE, align="center", space_after=0, line=1.05),
        ]))
    return els


def slide3():
    els = []
    els.append(image(os.path.join(ASSETS, "fig_arch.png"), 0.28, 1.42, 9.62, 4.58))
    # right column: implementation process
    els.append(text(10.05, 1.42, 3.0, 0.35,
                    [P("Implementation Process", 14, True, color=RED,
                       font=F_SERIF, space_after=0)]))
    steps = [
        ("Ingest & validate", "format · CRS · modality · co-registration · 512/64 tiling (GDAL)"),
        ("Plan & guard", "task id, tool selection from registry, parameter whitelist"),
        ("Execute", "specialists run per tile; every call writes to the trace"),
        ("Fuse & calibrate", "weighted geometric mean, High/Med/Low badges, disagreement rule"),
        ("Visualise", "change mask / boxes / caption streamed over the scene"),
        ("Report", "PDF/HTML audit pack + trace.json download"),
    ]
    y = 1.86
    for i, (t, d) in enumerate(steps):
        col = [BLUE, PURPLE, GREEN, TEAL, ORANGE, RED][i]
        els.append(oval(10.10, y, 0.34, 0.34, fill=col))
        els.append(text(10.10, y, 0.34, 0.34,
                        [P(str(i + 1), 11, True, color=WHITE, align="center",
                           space_after=0)], anchor="middle"))
        els.append(text(10.55, y - 0.05, 2.55, 0.75, [
            P(t, 10, True, color=INK, space_after=1, line=1.0),
            P(d, 7.6, False, color=GREY, space_after=0, line=1.05),
        ]))
        if i < 5:
            els.append(line(10.27, y + 0.36, 10.27, y + 0.72, color=col, lw=1.4))
        y += 0.72
    # tech stack strip
    sy = 6.22
    els.append(rect(0.28, sy, 12.77, 0.60, fill=WHITE, line=GREY, lw=1.0, radius=0.08))
    els.append(text(0.42, sy, 1.85, 0.60, [
        P("Technology stack", 10.5, True, color=DBLUE, font=F_SERIF,
          space_after=0, line=1.0),
        P("to be used", 10.5, True, color=DBLUE, font=F_SERIF, space_after=0, line=1.0),
    ], anchor="middle"))
    chips = ["Python", "FastAPI", "Streamlit", "GDAL / rasterio", "PyTorch",
             "Qwen2-VL (QLoRA)", "Grounding DINO", "UNet++ / BiT", "Docker"]
    cols = [BLUE, DBLUE, TEAL, GREEN, ORANGE, PURPLE, RED, ORANGE, BLUE]
    x = 2.45
    for i, c in enumerate(chips):
        w = 0.45 + len(c) * 0.062
        els.append(rect(x, sy + 0.13, w, 0.34, fill=cols[i], radius=0.17))
        els.append(text(x, sy + 0.13, w, 0.34,
                        [P(c, 7.6, True, color=WHITE, align="center", space_after=0)],
                        anchor="middle"))
        x += w + 0.10
    return els


def slide4():
    els = []
    boxes = [
        ("Feasibility", TINT_B, BLUE, [
            ("Proven open weights & open source — ", "Rs 0 licence cost."),
            ("Prototype already runs end-to-end ", "on a laptop (thin build, Window A)."),
            ("QLoRA 4-bit + 2B backbone ", "fits free Kaggle / Colab GPUs."),
            ("Tile-based 512 px / 64 overlap ", "→ sensor- & resolution-agnostic."),
        ]),
        ("Viability", TINT_G, GREEN, [
            ("15-day dated build plan ", "with daily checkpoints D1–D18."),
            ("3-model triage (M1/M4/M3) ", "matches pooled free compute exactly."),
            ("Demo-day safe: ", "deterministic fallback + cached-answer mode."),
            ("Metrics gates ", "on public splits before every submission."),
        ]),
        ("Practical Implementation", TINT_O, ORANGE, [
            ("Dockerised, offline deploy ", "for the evaluation venue."),
            ("Handles georeferenced Cartosat-2S + RISAT ", "pairs by design (P14)."),
            ("Normalised metric tables ", "(P15) in the final report."),
            ("Same tool contracts ", "from prototype → trained models."),
        ]),
    ]
    for i, (t, tint, col, items) in enumerate(boxes):
        x = 0.30 + i * 4.28
        els.append(rect(x, 1.42, 4.10, 2.02, fill=tint, line=col, lw=1.2, radius=0.07))
        paras = [P(t, 14, True, color=col, font=F_SERIF, space_after=4)]
        for b, rest in items:
            paras.append(rich([(b, 9.3, True, False, INK), (rest, 9.3, False, False, INK)],
                              space_after=3, line=1.05, bullet="•"))
        els.append(text(x + 0.16, 1.52, 3.80, 1.86, paras))
    # risk <-> strategy mirror panel
    py = 3.62
    els.append(rect(0.30, py, 12.73, 3.18, fill=WHITE, line=INK, lw=1.2, radius=0.08))
    els.append(oval(0.62, py + 0.92, 1.55, 1.34, fill=LGREY))
    els.append(text(0.62, py + 0.92, 1.55, 1.34, [
        P("Potential", 10.5, True, color=INK, align="center", space_after=0, line=1.05),
        P("Challenges", 10.5, True, color=INK, align="center", space_after=0, line=1.05),
        P("& Risks", 10.5, True, color=INK, align="center", space_after=0, line=1.05),
    ], anchor="middle"))
    els.append(oval(11.16, py + 0.92, 1.55, 1.34, fill="#DDEEDF"))
    els.append(text(11.16, py + 0.92, 1.55, 1.34, [
        P("Strategies", 10.5, True, color=GREEN, align="center", space_after=0, line=1.05),
        P("For Overcoming", 10.5, True, color=GREEN, align="center", space_after=0, line=1.05),
        P("Challenges", 10.5, True, color=GREEN, align="center", space_after=0, line=1.05),
    ], anchor="middle"))
    pairs = [
        ("Domain gap: ", "trained on Sentinel (10–20 m), evaluated on Cartosat-2S (~1 m) + RISAT.",
         "Tile-based, resolution-agnostic pipeline; ", "week-1 proxy tests on public VHR + L-band."),
        ("Tight 6 GB VRAM ", "at the evaluation venue.",
         "QLoRA 4-bit + 2B backbone; ", "heavy training on free 16 GB Kaggle GPUs."),
        ("Hidden eval annotations ", "& unpublished metric weights (G1).",
         "Train on the PS-named public benchmarks; ", "report normalised metric tables (P15)."),
        ("Demo-day failure risk ", "in front of judges.",
         "Guardrail fallback router + cached answers; ", "honest, scoped narration on stage."),
        ("15 days for 3 trained models.", "",
         "Binding triage M1/M4/M3; ", "dated checkpoints D1–D18; contracts frozen early."),
    ]
    y = py + 0.16
    for i, (rb, rr, sb, sr) in enumerate(pairs):
        n = f"{i+1:02d}"
        els.append(oval(2.42, y + 0.06, 0.30, 0.30, fill=RED))
        els.append(text(2.42, y + 0.06, 0.30, 0.30,
                        [P(n, 9, True, color=WHITE, align="center", space_after=0)],
                        anchor="middle"))
        els.append(text(2.84, y, 3.55, 0.56,
                        [rich([(rb, 9.2, True, False, INK), (rr, 9.2, False, False, INK)],
                              space_after=0, line=1.05)], anchor="middle"))
        els.append(oval(10.62, y + 0.06, 0.30, 0.30, fill=GREEN))
        els.append(text(10.62, y + 0.06, 0.30, 0.30,
                        [P(n, 9, True, color=WHITE, align="center", space_after=0)],
                        anchor="middle"))
        els.append(text(6.62, y, 3.88, 0.56,
                        [rich([(sb, 9.2, True, False, INK), (sr, 9.2, False, False, INK)],
                              space_after=0, line=1.05)], anchor="middle"))
        els.append(line(2.20, y + 0.21, 2.40, y + 0.21, color=RED, lw=1.2, arrow=True))
        els.append(line(11.14, y + 0.21, 10.94, y + 0.21, color=GREEN, lw=1.2, arrow=True))
        y += 0.60
    els.append(line(6.50, py + 0.15, 6.50, py + 3.03, color=LGREY, lw=1.0, dashed=True))
    return els


def slide5():
    els = []
    els.append(image(os.path.join(ASSETS, "fig_impact.png"), 0.30, 1.62, 7.30, 4.35))
    cards = [
        (2.30, 1.42, 3.60, 0.62, "ISRO / SAC analysts: ",
         "task answers with evidence in minutes, not scripting days."),
        (0.32, 2.52, 2.55, 1.05, "Disaster cells: ",
         "bi-temporal change maps + change-VQA for rapid damage triage."),
        (5.35, 2.52, 2.55, 1.05, "Agriculture dept: ",
         "water / crop queries in plain language — no GIS staff needed."),
        (0.32, 4.85, 2.55, 0.95, "Urban planners: ",
         "auditable built-up growth counts from masks, not eyeballing."),
        (5.35, 4.85, 2.55, 0.95, "Students & researchers: ",
         "open traces + reports make RS science reproducible."),
    ]
    for x, y, w, h, b, r in cards:
        els.append(rect(x, y, w, h, fill=CARDGREY, radius=0.06))
        els.append(text(x + 0.10, y + 0.06, w - 0.20, h - 0.12,
                        [rich([(b, 9.6, True, False, DBLUE), (r, 9.6, False, False, INK)],
                              space_after=0, line=1.08)], anchor="middle"))
    # benefits column
    els.append(rect(7.90, 1.42, 5.13, 0.42, fill=BLUE, radius=0.08))
    els.append(text(7.90, 1.42, 5.13, 0.42,
                    [P("Benefits of the solution", 14, True, color=WHITE,
                       font=F_SERIF, align="center", space_after=0)], anchor="middle"))
    bens = [
        ("Social", "Democratizes satellite intelligence — any officer queries in "
                   "plain language, the PS's own stated goal.", ORANGE),
        ("Economic", "One assistant replaces many single-task vendor pipelines; "
                     "saves analyst-hours on every tasking cycle.", GREEN),
        ("Environmental", "Faster spotting of floods, encroachment and water-body "
                          "change → quicker response on the ground.", TEAL),
    ]
    y = 2.02
    for t, b, col in bens:
        els.append(rect(7.90, y, 5.13, 1.28, fill=WHITE, line=col, lw=1.6, radius=0.09))
        els.append(text(8.06, y + 0.10, 1.35, 1.08,
                        [P(t, 12.5, True, color=col, font=F_SERIF, align="center",
                           space_after=0)], anchor="middle"))
        els.append(text(9.45, y + 0.10, 3.42, 1.08,
                        [P("“" + b + "”", 9.6, False, False, INK, space_after=0,
                           line=1.12)], anchor="middle"))
        y += 1.42
    els.append(rect(0.30, 6.30, 12.73, 0.50, fill=LGREY, radius=0.08))
    els.append(text(0.30, 6.30, 12.73, 0.50,
                    [P("“Every answer with evidence, confidence and an audit trail — "
                       "satellite intelligence anyone can query.”", 12, True, False,
                       DBLUE, font=F_SERIF, align="center", space_after=0)],
                    anchor="middle"))
    return els


def slide6():
    els = []
    els.append(rect(0.30, 1.42, 12.73, 5.38, fill=WHITE, line=INK, lw=1.2, radius=0.08))
    L = [
        P("Official problem statement & portals", 14, True, color=PURPLE,
          font=F_SERIF, space_after=4),
        rich([("SIH 2026 Problem Statement 26167 (ISRO): ", 10, True, False, INK),
              ("“SatQuery AI — An Interactive Vision-Language Assistant for "
               "Multimodal Remote Sensing Image Analysis through Text Queries” · ",
               10, False, False, INK),
              ("sih.gov.in", 10, False, False, BLUE)], space_after=3, line=1.08, bullet="•"),
        rich([("ISRO / Space Applications Centre: ", 10, True, False, INK),
              ("evaluation set of pre-georeferenced Cartosat-2S + RISAT pairs (P14) · ",
               10, False, False, INK),
              ("isro.gov.in", 10, False, False, BLUE)], space_after=8, line=1.08, bullet="•"),
        P("Datasets & benchmarks named by the PS", 14, True, color=PURPLE,
          font=F_SERIF, space_after=4),
        rich([("BigEarthNet.txt ", 10, True, False, INK),
              ("(primary adaptation data — 9.55 M annotations over 464 K S1–S2 pairs): "
               "arXiv:2603.29630 · huggingface.co/datasets/BIFOLD-BigEarthNetv2-0/"
               "BigEarthNet.txt", 10, False, False, INK)],
             space_after=3, line=1.08, bullet="•"),
        rich([("VRSBench ", 10, True, False, INK),
              ("(NeurIPS'24 D&B — captioning, grounding, VQA): arXiv:2406.12384 · "
               "github.com/lx709/VRSBench", 10, False, False, INK)],
             space_after=3, line=1.08, bullet="•"),
        rich([("RSVQA ", 10, True, False, INK),
              ("(Lobry et al., IEEE TGRS 2020) — single-image VQA volume.",
               10, False, False, INK)], space_after=3, line=1.08, bullet="•"),
        rich([("CDVQA ", 10, True, False, INK),
              ("(Yuan et al., TGRS 2022 — bi-temporal change VQA + semantic change "
               "maps): arXiv:2112.06343 · github.com/YZHJessica/CDVQA",
               10, False, False, INK)], space_after=3, line=1.08, bullet="•"),
        rich([("TAMMI (2025) ", 10, True, False, INK),
              ("— VQA over co-located VHR RGB + multispectral + SAR triplets; "
               "LEVIR-CD / WHU-CD for change-mask volume "
               "(github.com/justchenhao/Levir-CD).", 10, False, False, INK)],
             space_after=3, line=1.08, bullet="•"),
    ]
    R_ = [
        P("Models & methods", 14, True, color=PURPLE, font=F_SERIF, space_after=4),
        rich([("Qwen2-VL ", 10, True, False, INK),
              ("(2 B backbone, QLoRA 4-bit adaptation): arXiv:2409.12191 · QLoRA "
               "arXiv:2305.14314 · LoRA arXiv:2106.09685", 10, False, False, INK)],
             space_after=3, line=1.08, bullet="•"),
        rich([("Grounding DINO ", 10, True, False, INK),
              ("(text-guided grounding base): arXiv:2303.05499",
               10, False, False, INK)], space_after=3, line=1.08, bullet="•"),
        rich([("UNet++ / BiT ", 10, True, False, INK),
              ("(change & patch-classifier heads): arXiv:1807.10165",
               10, False, False, INK)], space_after=8, line=1.08, bullet="•"),
        P("Our groundwork (this submission)", 14, True, color=PURPLE,
          font=F_SERIF, space_after=4),
        rich([("Clause-level PS analysis: ", 10, True, False, INK),
              ("all 15 binding clauses (P1–P15) mapped to design decisions; "
               "5 known PS gaps logged (G1–G5).", 10, False, False, INK)],
             space_after=3, line=1.08, bullet="•"),
        rich([("Working thin prototype: ", 10, True, False, INK),
              ("real ingestion (GDAL), registry + guardrail + fallback router, "
               "trace store, overlays and PDF reports — screenshots on slide 2.",
               10, False, False, INK)], space_after=3, line=1.08, bullet="•"),
        rich([("Dated 15-day build plan: ", 10, True, False, INK),
              ("three trained models (M1/M4/M3) with daily checkpoints and "
               "benchmark gates for the 20 Sep final submission.",
               10, False, False, INK)], space_after=3, line=1.08, bullet="•"),
    ]
    gates = [
        ("Captioning", "VRSBench test", "CIDEr · BLEU-4 · METEOR · ROUGE-L", BLUE),
        ("Grounding", "VRSBench refs", "box IoU ≥ 0.5 acc · COCO AP", ORANGE),
        ("Single-image VQA", "RSVQA + VRSBench", "exact match · soft accuracy", GREEN),
        ("Change-VQA", "CDVQA test", "exact match over answer categories", PURPLE),
        ("Change map", "CDVQA / QAG-360K", "IoU · F1 · mIoU", TEAL),
    ]
    els.append(text(0.62, 4.92, 12.1, 0.35,
                    [P("Benchmark gates for the 20 Sep submission — zero-shot → "
                       "adapted delta is our P12 compliance proof", 12.5, True,
                       color=PURPLE, font=F_SERIF, space_after=0)]))
    gw = 2.38
    for i, (t, d, m, col) in enumerate(gates):
        x = 0.62 + i * (gw + 0.08)
        els.append(rect(x, 5.32, gw, 0.92, fill=WHITE, line=col, lw=1.4, radius=0.07))
        els.append(text(x + 0.08, 5.38, gw - 0.16, 0.82, [
            P(t, 9.2, True, color=col, align="center", space_after=1, line=1.0),
            P(d, 7.8, False, color=GREY, align="center", space_after=1, line=1.0),
            P(m, 7.6, True, color=INK, align="center", space_after=0, line=1.05),
        ]))
    els.append(text(0.62, 6.32, 12.1, 0.3,
                    [P("All scores normalised before combining (P15); per-tool "
                       "latency and confidence recorded in every trace.", 8.6,
                       False, True, GREY, align="center", space_after=0)]))
    els.append(text(0.62, 1.62, 6.05, 5.0, L))
    els.append(text(6.95, 1.62, 5.80, 5.0, R_))
    els.append(line(6.72, 1.62, 6.72, 4.75, color=LGREY, lw=1.0, dashed=True))
    return els


# ---------------------------------------------------------------- PPTX text edits
def _apply_slide_texts(prs, slide, idx):
    if idx == 0:
        _set_ph(slide, "Title 7", "SMART INDIA HACKATHON 2026", size=40, bold=True,
                color=DBLUE, font=F_GAR)
        _set_ph(slide, "Subtitle 3",
                "SatQuery AI — ask your satellite in plain language",
                size=17, bold=False, italic=True, color=GREY, font=F_SERIF)
        _set_ph(slide, "TextBox 9",
                "Problem Statement ID – 26167\n"
                "Problem Statement Title – SatQuery AI: An Interactive Vision-Language\n"
                "Assistant for Multimodal Remote Sensing Image Analysis through Text Queries\n"
                "Theme – Space Technology\n"
                "PS Category – Software\n"
                "Team ID – [ to be filled ]\n"
                "Team Name – x64",
                size=18, color=INK, font=F_ARIAL)
        # re-style: bold labels
        for sh in slide.shapes:
            if sh.name == "TextBox 9":
                sh.left = Inches(0.40); sh.top = Inches(2.20)
                sh.width = Inches(6.60); sh.height = Inches(4.95)
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if "–" in r.text:
                            head, _, tail = r.text.partition("–")
                            r.text = head
                            r.font.bold = True
                            r2 = p.add_run(); r2.text = "– " + tail
                            r2.font.size = Pt(18); r2.font.name = F_ARIAL
                            r2.font.color.rgb = RGBColor.from_string(
                                GREY if "filled" in tail else INK)
        return
    titles = {1: "SATQUERY AI", 2: "TECHNICAL APPROACH", 3: "FEASIBILITY AND VIABILITY",
              4: "IMPACT AND BENEFITS", 5: "RESEARCH AND REFERENCES"}
    _set_ph(slide, "Title 1", titles[idx], size=30, bold=True, color=INK, font=F_SERIF,
            align="center")
    for sh in slide.shapes:
        if sh.name == "Title 1":
            sh.left = Inches(1.90); sh.top = Inches(0.10)
            sh.width = Inches(9.40); sh.height = Inches(0.85)
    ptr = {1: "Proposed solution (idea / solution / prototype) · how it addresses the "
              "problem · innovation & uniqueness",
           2: "Technologies to be used · methodology & process for implementation "
              "(flow charts / images / working prototype)",
           3: "Feasibility analysis · potential challenges & risks · strategies for "
              "overcoming these challenges",
           4: "Potential impact on the target audience · benefits of the solution "
              "(social, economic, environmental)",
           5: "Details / links of the reference and research work"}[idx]
    tb = slide.shapes.add_textbox(Inches(1.90), Inches(1.02), Inches(9.40), Inches(0.30))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ptr
    r.font.size = Pt(9.5); r.font.italic = True; r.font.name = F_ARIAL
    r.font.color.rgb = RGBColor.from_string(GREY)
    # footer + number
    for sh in slide.shapes:
        if sh.name == "Footer Placeholder 6":
            tf = sh.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "SMART INDIA HACKATHON 2026"
            r.font.size = Pt(13); r.font.bold = True; r.font.name = F_SERIF
            r.font.color.rgb = RGBColor.from_string(WHITE)
        if sh.name == "Slide Number Placeholder 5":
            tf = sh.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(idx + 1)
            r.font.size = Pt(13); r.font.bold = True; r.font.name = F_SERIF
            r.font.color.rgb = RGBColor.from_string(WHITE)
        if sh.name.startswith("Oval") and abs(Emu(sh.left).inches - 0.36) < 0.05 \
                and abs(Emu(sh.top).inches - 0.28) < 0.05:
            tf = sh.text_frame; tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "x64"
            r.font.size = Pt(15); r.font.bold = True; r.font.name = F_ARIAL
            r.font.color.rgb = RGBColor.from_string(PURPLE)
            sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(WHITE)
            sh.line.color.rgb = RGBColor.from_string(PURPLE)


def main():
    # extract template logo for the PDF renderer
    logo = os.path.join(ASSETS, "sih_logo.png")
    if not os.path.exists(logo):
        import zipfile
        z = zipfile.ZipFile(TEMPLATE)
        open(logo, "wb").write(z.read("ppt/media/image2.png"))
        import io
        from PIL import Image
        im = Image.open(io.BytesIO(z.read("ppt/media/image1.png")))
        w, h = im.size
        im.crop((0, 0, int(w * 0.40084), h)).save(
            os.path.join(ASSETS, "sih_bulb.png"))
    # downsample raster assets to sane print resolution (keeps the PDF small)
    from PIL import Image
    limits = {"sih_logo.png": 1000, "sih_bulb.png": 1200,
              "overlay_change.png": 512, "overlay_ground.png": 512,
              "overlay_caption.png": 512}
    for fn, maxw in limits.items():
        p = os.path.join(ASSETS, fn)
        im = Image.open(p)
        if im.width > maxw:
            im = im.convert("RGBA").resize(
                (maxw, round(im.height * maxw / im.width)), Image.LANCZOS)
            im.save(p)
    specs = [slide1(), slide2(), slide3(), slide4(), slide5(), slide6()]
    ptrs = [None] + [c["fg"] for c in []]  # placeholder
    chrome = []
    titles = ["", "SATQUERY AI", "TECHNICAL APPROACH", "FEASIBILITY AND VIABILITY",
              "IMPACT AND BENEFITS", "RESEARCH AND REFERENCES"]
    pointers = [None,
                "Proposed solution (idea / solution / prototype) · how it addresses "
                "the problem · innovation & uniqueness",
                "Technologies to be used · methodology & process for implementation "
                "(flow charts / images / working prototype)",
                "Feasibility analysis · potential challenges & risks · strategies for "
                "overcoming these challenges",
                "Potential impact on the target audience · benefits of the solution "
                "(social, economic, environmental)",
                "Details / links of the reference and research work"]
    for i in range(6):
        chrome.append(chrome_for(i, titles[i], pointers[i], i + 1))
    render_pptx(specs, OUT_PPTX)
    print("pptx ->", OUT_PPTX)
    render_pdf(specs, chrome, OUT_PDF)
    print("pdf  ->", OUT_PDF)


if __name__ == "__main__":
    main()
